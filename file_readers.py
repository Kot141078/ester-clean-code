# -*- coding: utf-8 -*-
"""
file_readers.py — unifitsirovannye ridery faylov.
Podderzhka: txt, md, html, pdf, docx, xlsx, pptx, json, py, fb2, epub.
"""

import io
import os
import re
import json
import logging
import xml.etree.ElementTree as ET
from modules.memory.facade import memory_add, ESTER_MEM_FACADE

# Biblioteki (import vnutri funktsiy, bezopasnyy)
try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import docx
except ImportError:
    docx = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    import ebooklib
    from ebooklib import epub
except ImportError:
    ebooklib = None

# Regular rule for determining the language
_LANG_RX = re.compile(r"[A-Yaa-yaEe]")

def detect_and_read(filename: str, raw_data: bytes) -> tuple[list, str]:
    """Main entry point."""
    ext = os.path.splitext(filename)[1].lower()
    
    if ext == ".pdf":
        return _read_pdf(raw_data)
    elif ext in [".docx", ".doc"]:
        return _read_docx(raw_data)
    elif ext in [".pptx", ".ppt"]:
        return _read_pptx(raw_data)
    elif ext in [".xlsx", ".xls", ".csv", ".tsv"]:
        return _read_spreadsheet(raw_data, ext)
    elif ext in [".json", ".jsonl"]:
        return _read_json(raw_data)
    elif ext == ".fb2":
        return _read_fb2(raw_data)
    elif ext == ".epub":
        return _read_epub(raw_data)
    else:
        return _read_text_generic(raw_data)

def _read_text_generic(raw: bytes):
    try:
        text = raw.decode("utf-8", errors="ignore")
    except:
        text = str(raw)
    return [{"text": text}], text

def _read_pdf(raw: bytes):
    text = ""
    if not pypdf:
        return [{"text": "[ERROR: pypdf not installed]"}], ""
    try:
        with io.BytesIO(raw) as f:
            reader = pypdf.PdfReader(f)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
    except Exception as e:
        text = f"[PDF Error: {e}]"
    return [{"text": text}], text

def _read_docx(raw: bytes):
    text = ""
    if not docx:
        return [{"text": "[ERROR: python-docx not installed]"}], ""
    try:
        with io.BytesIO(raw) as f:
            doc = docx.Document(f)
            text = "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        text = f"[DOCX Error: {e}]"
    return [{"text": text}], text

def _read_pptx(raw: bytes):
    text = ""
    if not pptx:
        return [{"text": "[ERROR: python-pptx not installed]"}], ""
    try:
        with io.BytesIO(raw) as f:
            prs = pptx.Presentation(f)
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            text = "\n".join(full_text)
    except Exception as e:
        text = f"[PPTX Error: {e}]"
    return [{"text": text}], text

def _read_spreadsheet(raw: bytes, ext: str):
    text = ""
    if not pd:
        return [{"text": "[ERROR: pandas not installed]"}], ""
    try:
        with io.BytesIO(raw) as f:
            if ext == ".csv":
                df = pd.read_csv(f)
            elif ext == ".tsv":
                df = pd.read_csv(f, sep="\t")
            else:
                df = pd.read_excel(f)
            text = df.head(500).to_string(index=False)
    except Exception as e:
        text = f"[Spreadsheet Error: {e}]"
    return [{"text": text}], text

def _read_json(raw: bytes):
    try:
        data = json.loads(raw.decode("utf-8", errors="ignore"))
        text = json.dumps(data, ensure_ascii=False, indent=2)
    except Exception as e:
        text = f"[JSON Error: {e}]"
    return [{"text": text}], text

def _read_fb2(raw: bytes):
    text = ""
    try:
        # FB2 is CML. Let's try to parse with the standard ElementTree
        # First we decode, FB2 is often in UTF-8 or Windows-1251
        try:
            xml_content = raw.decode("utf-8")
        except:
            xml_content = raw.decode("windows-1251", errors="ignore")
            
        # We remove the space if it gets in the way (simple cleaning)
        xml_content = re.sub(r' xmlns="[^"]+"', '', xml_content, count=1)
        
        root = ET.fromstring(xml_content)
        
        # Izvlekaem tekst iz body
        body = root.find("body")
        if body is None:
            # If water is not found directly, we search recursively
            body = root.find(".//body")
            
        if body is not None:
            # Sobiraem vse paragrafy (p) i zagolovki (title)
            chunks = []
            for elem in body.iter():
                if elem.tag in ['p', 'v', 'subtitle', 'title'] and elem.text:
                    chunks.append(elem.text.strip())
            text = "\n".join(chunks)
        else:
            text = "[FB2 Error: Body not found]"
            
    except Exception as e:
        text = f"[FB2 Error: {e}]"
    return [{"text": text}], text

def _read_epub(raw: bytes):
    text = ""
    if not ebooklib:
        return [{"text": "[ERROR: ebooklib not installed. pip install EbookLib]"}], ""
    try:
        # EbookLib requires a file on disk, but can we cheat through Bitsio?
        # Unfortunately, EbookLib does not work well with BeatIO directly.
        # Therefore, save it to a temporary file
        import tempfile
        with tempfile.NamedTemporaryFile(delete=False, suffix=".epub") as tmp:
            tmp.write(raw)
            tmp_path = tmp.name
        
        book = epub.read_epub(tmp_path)
        
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                # This is NTML, you need to clean the tags. Simple regex
                html_content = item.get_content().decode("utf-8", errors="ignore")
                clean_text = re.sub('<[^<]+?>', '', html_content)
                text += clean_text + "\n"
        
        try:
            os.remove(tmp_path)
        except: pass
        
    except Exception as e:
        text = f"[EPUB Error: {e}]"
    return [{"text": text}], text